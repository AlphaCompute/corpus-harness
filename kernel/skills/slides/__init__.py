"""A deck from the same markdown a report is written in, because the two are usually the
same material and writing it twice is how they come to disagree."""

import re
from pathlib import Path

# A picture on a line of its own. Inside a sentence it stays prose, because a slide has
# nowhere to put an image mid-paragraph.
_IMAGE = re.compile(r"^!\[[^\]]*\]\(([^)]+)\)\s*$")

# What a bullet on a slide has no way of showing. The markers are stripped rather than
# rendered, so `**quarter**` reads as a word instead of as four asterisks and a word.
_EMPHASIS = re.compile(r"\*\*(.+?)\*\*|\*([^*\n]+)\*|`([^`]+)`")

# Enough of the slide for a table or a picture, under a title placeholder that ends around
# an inch and a half down.
_TOP = 1.8
_MARGIN = 0.7


def deck(content, path="deck.pptx", title=None, template=None):
    """Writes a .pptx and returns where it landed.

    `content` is text in the markdown subset `documents.blocks` understands. A `#` heading
    opens the deck's title slide, every deeper heading starts a slide of its own, and the
    lists, paragraphs, tables and images under it are what goes on it:

        slides.deck("# Q3\\n\\n## Revenue\\n\\n- up 12%\\n\\n![](chart.png)", "q3.pptx")

    `template` is a .pptx to inherit from — the corporate one, usually — and decides the
    slide size, the fonts and the colours. Without one you get PowerPoint's own 4:3
    default, which is what python-pptx ships.
    """
    from pptx import Presentation

    show = Presentation(template) if template else Presentation()
    for head, level, blocks in _slides(content):
        head = head or title or Path(path).stem
        prose = [block for block in blocks if block[0] == "text"]
        if level == 1:
            _opening(show, head, prose)
        elif prose:
            _bullets(show, head, prose)
        # A table and a picture each get a slide of their own, under the same title: they
        # are the size of the slide, and a body placeholder is already holding the prose.
        for kind, payload in blocks:
            if kind == "table":
                _table(show, head, payload)
            elif kind == "image":
                _picture(show, head, payload)
    if not show.slides:
        _opening(show, title or Path(path).stem, [])
    show.save(str(path))
    return str(path)


def _slides(text):
    """The deck as `(title, heading level, blocks)`, where a block is `("text", line)`,
    `("table", rows)` or `("image", path)`.

    The markdown is read with `documents`' own dispatch rather than a second copy of it.
    Those names are private to that module, and this reaches past that on purpose: a deck
    that understood a different subset than the report it was cut from would differ from it
    in exactly the places nobody thinks to check.
    """
    from documents import _HEADING, _RULE, _item, _kind, _take

    slides = []
    lines = text.splitlines()
    at = 0
    while at < len(lines):
        line = lines[at]
        if not line.strip():
            at += 1
            continue

        picture = _IMAGE.match(line)
        if picture is not None:
            _current(slides).append(("image", picture.group(1)))
            at += 1
            continue

        kind = _kind(line)
        if kind == "heading":
            heading = _HEADING.match(line)
            slides.append((_plain(heading.group(2)), len(heading.group(1)), []))
            at += 1
        elif kind == "row":
            rows, at = _take(lines, at, lambda line: _kind(line) == "row")
            cells = [
                [cell.strip() for cell in row.strip().strip("|").split("|")]
                for row in rows
                if not _RULE.match(row)
            ]
            _current(slides).append(("table", cells))
        elif kind == "item":
            raw, at = _take(lines, at, lambda line: _kind(line) == "item")
            for item in raw:
                _current(slides).append(("text", _plain(_item(item))))
        else:
            raw, at = _take(lines, at, lambda line: line.strip() and _kind(line) is None)
            # A paragraph's line breaks are the width of whoever typed it; a slide's are
            # the width of the slide.
            _current(slides).append(("text", _plain(" ".join(piece.strip() for piece in raw))))
    return slides


def _current(slides):
    """The blocks of the slide being filled, opening an untitled one for material that
    arrived before any heading did."""
    if not slides:
        slides.append((None, 2, []))
    return slides[-1][2]


def _plain(text):
    return _EMPHASIS.sub(lambda match: next(filter(None, match.groups())), text)


def _slide(show, layout, title):
    slide = show.slides.add_slide(show.slide_layouts[layout])
    slide.shapes.title.text = title
    return slide


def _opening(show, title, prose):
    slide = _slide(show, 0, title)
    # The subtitle placeholder is the layout's second, and a template that dropped it
    # leaves the opening line to the slide that follows rather than raising here.
    body = [holder for holder in slide.placeholders if holder.placeholder_format.idx == 1]
    if prose and body:
        body[0].text = "\n".join(text for _, text in prose)


def _bullets(show, title, prose):
    slide = _slide(show, 1, title)
    frame = slide.placeholders[1].text_frame
    frame.text = prose[0][1]
    for _, text in prose[1:]:
        frame.add_paragraph().text = text


def _table(show, title, rows):
    from pptx.util import Inches

    columns = max(len(row) for row in rows)
    width = show.slide_width - Inches(2 * _MARGIN)
    height = show.slide_height - Inches(_TOP + _MARGIN)
    shape = _slide(show, 5, title).shapes.add_table(
        len(rows), columns, Inches(_MARGIN), Inches(_TOP), width, height
    )
    for down, row in enumerate(rows):
        for across in range(columns):
            # A row short of the rectangle is the writer's slip, and an empty cell says so
            # more usefully than an exception thrown while saving.
            shape.table.cell(down, across).text = _plain(
                row[across] if across < len(row) else ""
            )


def _picture(show, title, source):
    from pptx.util import Inches

    slide = _slide(show, 5, title)
    room = (show.slide_width - Inches(2 * _MARGIN), show.slide_height - Inches(_TOP + _MARGIN))
    picture = slide.shapes.add_picture(source, Inches(_MARGIN), Inches(_TOP), width=room[0])
    # Added to the width first because that is the dimension a chart usually wants; a tall
    # image overflows the slide bottom silently, so it is measured and re-fitted by height.
    if picture.height > room[1]:
        picture.width = int(picture.width * room[1] / picture.height)
        picture.height = room[1]
    picture.left = int((show.slide_width - picture.width) / 2)
